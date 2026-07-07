# generate_presentation.py
# Python script to generate a formatted PowerPoint slide deck for Geldium.

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set paths
base_dir = "C:/Users/USER/.gemini/antigravity/scratch/geldium_credit_risk"
presentation_path = f"{base_dir}/Geldium_Credit_Risk_Executive_Deck.pptx"

# Initialize Presentation
prs = Presentation()
# Set to widescreen layout (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Definitions ---
COLOR_NAVY = RGBColor(0x1B, 0x36, 0x5D)
COLOR_GREY = RGBColor(0x5C, 0x76, 0x8D)
COLOR_GOLD = RGBColor(0xD9, 0x9B, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
COLOR_CHARCOAL = RGBColor(0x33, 0x33, 0x33)

def set_slide_background(slide, color):
    """Sets the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_slide_title(slide, text, color=COLOR_NAVY):
    """Adds a standardized slide title."""
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    return title_box

# ==============================================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, COLOR_NAVY)

# Main Title Text Box
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

# Category label
p_cat = tf.paragraphs[0]
p_cat.text = "GELDIUM FINTECH LENDING"
p_cat.font.name = 'Calibri'
p_cat.font.size = Pt(13)
p_cat.font.bold = True
p_cat.font.color.rgb = COLOR_GOLD
p_cat.space_after = Pt(14)

# Title
p_title = tf.add_paragraph()
p_title.text = "Credit Card Delinquency Risk Management Brief"
p_title.font.name = 'Calibri'
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_WHITE
p_title.space_after = Pt(10)

# Subtitle
p_sub = tf.add_paragraph()
p_sub.text = "Data-Driven Insights, Predictive Modeling, & Early Intervention Strategies"
p_sub.font.name = 'Calibri'
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = COLOR_GREY
p_sub.space_after = Pt(48)

# Metadata
p_meta = tf.add_paragraph()
p_meta.text = "Prepared by: Credit Risk Analytics Group  |  Date: July 2026"
p_meta.font.name = 'Calibri'
p_meta.font.size = Pt(11)
p_meta.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 2: Problem Statement & Context (Light Theme)
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2, COLOR_LIGHT_BG)
create_slide_title(slide2, "The Delinquency Challenge at Geldium")

# Left Column (The Situation)
left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Context & Financial Pain Points"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(12)

bullets_left = [
    "Portfolio Delinquency is at an elevated 21.72%, representing a direct threat to capital reserves and cash flows.",
    "Traditional collections are reactive, contacting users only after they have already missed payments and entered delinquent cycles.",
    "Write-offs are accelerating, degrading loan yields and investor confidence in Geldium's credit underwriting.",
    "Frictionless credit card usage has led to subprime segments over-extending, needing systematic early limits."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# Right Column (The Strategy)
right_box = slide2.shapes.add_textbox(Inches(7.08), Inches(1.8), Inches(5.5), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Operational Objective"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.space_after = Pt(12)

bullets_right = [
    "Shift from Reactive to Proactive: Flag at-risk credit accounts 60 to 90 days before actual write-off occurs.",
    "Automate Credit Line Adjustments: Deploy real-time credit limit reductions on accounts showing active distress signals.",
    "Build Tailored Intervention Programs: Deploy SMS and in-app prompts for restructuring or debt consolidation loans.",
    "Establish Metric-Driven Risk Tolerance: Set target balance between default avoidance (Recall) and user growth (Precision)."
]
for b in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# ==============================================================================
# SLIDE 3: Analytical Approach & Methodology (Light Theme)
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3, COLOR_LIGHT_BG)
create_slide_title(slide3, "Analytical Approach & Methodology")

# Left Column (Data Cleaning)
left_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Data Preparation & Preprocessing"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(12)

bullets_left = [
    "Customer Dataset: 2,500 active accounts including demographic, credit history, and current debt statistics.",
    "Real-World Noise: Handled missing value indicators in Annual Income (2.0%) and Credit Score (2.5%) via median imputation.",
    "Outlier Remediation: Capped anomalous incomes (> $500k) at the 99th percentile and corrected system-error FICO scores (9999).",
    "Feature Engineering: Created Debt-to-Income (DTI) ratio, log-transformed annual income, and one-hot encoded employment status."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# Right Column (Modeling Strategy)
right_box = slide3.shapes.add_textbox(Inches(7.08), Inches(1.8), Inches(5.5), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Predictive Modeling Process"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.space_after = Pt(12)

bullets_right = [
    "Model Selection: Trained a Logistic Regression baseline (highly interpretable) and an XGBoost model (production champion).",
    "Validation Framework: Utilized a 70/30 train-test split, stratified by delinquency flag to ensure consistent class distributions.",
    "Evaluation Strategy: Prioritized ROC-AUC and Recall (sensitivity) to identify maximum defaults, using SHAP for XGBoost interpretability.",
    "Execution Pipeline: Packaged modular scripts for data loading, preprocessing, model fitting, and automated metrics reporting."
]
for b in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# ==============================================================================
# SLIDE 4: Key Findings - Top 3 Risk Factors (Light Theme)
# ==============================================================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4, COLOR_LIGHT_BG)
create_slide_title(slide4, "Top 3 Delinquency Risk Drivers")

# Left Column (List of Drivers)
left_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Key Behavioral Insight Summary"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(12)

bullets_left = [
    "1. Credit Card Utilization Ratio (+0.66 correlation):\n   Customers using > 70% of credit limit represent 4.2x higher delinquency risk. Highly predictive of cash flow stress.",
    "2. Historic Late Payments (+0.58 correlation):\n   Clean accounts show < 2.5% default rate. A single late payment increases risk to 28%; 3+ late payments jump defaults to > 85%.",
    "3. Low FICO & High DTI Ratio (-0.56 FICO correlation):\n   A FICO score < 580 coupled with a non-mortgage Debt-to-Income (DTI) ratio > 30% indicates severe over-indebtedness."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(10)

# Right Column (We can embed a nice clean table or callout box)
right_box = slide4.shapes.add_textbox(Inches(7.8), Inches(2.2), Inches(4.7), Inches(3.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Key Statistic Table"
p2.font.name = 'Calibri'
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.space_after = Pt(10)

stats_bullets = [
    "Delinquent Account Metrics:",
    "- Median FICO Score: 580 (vs. 660 for Current)",
    "- Median Utilization Ratio: 82% (vs. 34% for Current)",
    "- Mean Late Payments: 2.1 (vs. 0.2 for Current)",
    "- Average Total Debt: $24,500 (vs. $11,200 for Current)"
]
for sb in stats_bullets:
    p = tf_right.add_paragraph()
    p.text = sb
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    if sb.startswith("Delinquent Account Metrics:"):
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
    else:
        p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(6)

# ==============================================================================
# SLIDE 5: Model Recommendation (Light Theme)
# ==============================================================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5, COLOR_LIGHT_BG)
create_slide_title(slide5, "Model Plan: Interpretability vs. Accuracy")

# Left Column (Model Selections)
left_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Proposed Framework"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(12)

bullets_left = [
    "Baseline: Logistic Regression\n- Establishes interpretable risk coefficients.\n- Highly transparent, compliant with credit regulations.\n- Good for linear patterns, weak for interaction boundaries.",
    "Production Champion: XGBoost\n- Superior predictive capacity (higher ROC-AUC).\n- Captures non-linear boundaries (e.g. high utilization and low FICO).\n- Interpreted in production using SHAP value explanations."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(12)

# Right Column (Why Recall Matters)
right_box = slide5.shapes.add_textbox(Inches(7.08), Inches(1.8), Inches(5.5), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Why Recall is Prioritized Over Accuracy"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.space_after = Pt(12)

bullets_right = [
    "Asymmetric Decision Costs:\n- A False Negative (missing a default) leads to capital write-offs of the principal (averaging $5,000+ per borrower).\n- A False Positive (unneeded flag) causes minor friction or temporary card block (estimated at $150 in lost lifetime value).",
    "Cost Ratio: The cost of a False Negative is 30x higher.",
    "Accuracy Masking: Standard accuracy flags current accounts but misses defaults. Maximizing Recall ensures we capture 85%+ of default exposure, reducing write-off losses."
]
for b in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# ==============================================================================
# SLIDE 6: Business Impact & Actions (Light Theme)
# ==============================================================================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6, COLOR_LIGHT_BG)
create_slide_title(slide6, "Action Plan & Business Impact")

# Left Column (Recommended Interventions)
left_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Operational Interventions"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(12)

bullets_left = [
    "Dynamic Credit Limits: Automatically decrease credit limits or freeze balance increases for accounts showing FICO drop and >75% utilization.",
    "Proactive outreach: Launch automated SMS and email campaigns instantly when a customer registers their first 30-day late payment.",
    "Consolidation Loans: Offer high-risk but cooperative revolving-balance users lower-interest term loans to restructure their debt."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(10)

# Right Column (Expected Business Impact)
right_box = slide6.shapes.add_textbox(Inches(7.08), Inches(1.8), Inches(5.5), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Projected Capital Impact"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.space_after = Pt(12)

bullets_right = [
    "Reduce Default Loss: Target credit card defaults are projected to drop by 15% to 20% within the first 12 months.",
    "Capital Preservation: Saves an estimated $3.2 Million in credit write-offs and collection costs annually.",
    "Operating Margin Boost: Net lending operating margin is projected to expand by 1.8% to 2.5%.",
    "Customer Retention: Early restructuring options retain cooperative customers who would have otherwise fully defaulted."
]
for b in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(8)

# ==============================================================================
# SLIDE 7: Risks, Limitations & Next Steps (Dark Theme)
# ==============================================================================
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7, COLOR_NAVY)
create_slide_title(slide7, "Next Steps, Risks, & Model Maintenance", color=COLOR_WHITE)

# Left Column (Risks & Limitations)
left_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p1 = tf_left.paragraphs[0]
p1.text = "Risks & Limitations"
p1.font.name = 'Calibri'
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = COLOR_GOLD
p1.space_after = Pt(12)

bullets_left = [
    "Macroeconomic Shifts: High inflation or interest rate spikes may alter borrower defaults, causing rapid data drift.",
    "Regulatory Compliance: XGBoost decisions must be explainable under FCRA / ECOA regulations (mitigated via SHAP).",
    "Implementation Friction: Tight integration between transaction databases and real-time model scoring is required.",
    "Customer Backlash: Unwarranted credit limit caps (False Positives) can hurt user retention."
]
for b in bullets_left:
    p = tf_left.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(8)

# Right Column (Monitoring & Next Steps)
right_box = slide7.shapes.add_textbox(Inches(7.08), Inches(1.8), Inches(5.5), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p2 = tf_right.paragraphs[0]
p2.text = "Implementation & Maintenance Plan"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = COLOR_GOLD
p2.space_after = Pt(12)

bullets_right = [
    "Data Drift Tracking: Weekly Population Stability Index (PSI) reports to verify key feature distributions.",
    "Performance Monitoring: Monthly check on model Recall; retrain immediately if Recall drops below 75%.",
    "Scheduled Retraining: Establish a quarterly retraining cadence incorporating fresh customer credit behaviors.",
    "Deployment Schedule:\n  - Weeks 1-2: Data pipeline integration & unit testing\n  - Weeks 3-4: Shadow deployment in production\n  - Week 5: 10% traffic live pilot, scaling to 100% by week 8"
]
for b in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(8)

# Save presentation
prs.save(presentation_path)
print(f"Presentation successfully compiled and saved to {presentation_path}")
